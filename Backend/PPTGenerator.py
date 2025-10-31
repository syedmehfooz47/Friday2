# -*- coding: utf-8 -*-
"""
PowerPoint Generator - Creates professional presentations using Groq AI
Generates well-designed, multi-slide PPTX files
"""

import os
import time
from pathlib import Path
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from .logger import Logger
from .llm_handler import llm_handler

load_dotenv()


class PPTGenerator:
    """Generate professional PowerPoint presentations"""
    
    def __init__(self):
        self.data_folder = Path(__file__).parent.parent / "Data" / "GeneratedDocuments"
        self.data_folder.mkdir(parents=True, exist_ok=True)
    
    def generate_ppt(self, topic: str, slides: int = 7) -> tuple[str, str]:
        """
        Generate a professional PowerPoint presentation
        
        Args:
            topic: Topic for the presentation
            slides: Number of slides to generate (minimum 7)
            
        Returns:
            Tuple of (response message, file path)
        """
        username = os.getenv("Username", "Boss")
        
        if slides < 1:
            return f"Invalid slide count, Boss. Please specify at least 1 slide.", None
        
        Logger.log(f"Generating {slides}-slide presentation on topic: '{topic}'", "PPT")
        
        try:
            # Generate content using Groq
            content_prompt = f"""You are a professional presentation designer. Create a comprehensive presentation about "{topic}" with exactly {slides} slides.

For each slide, provide:
1. Slide title (clear and concise)
2. Bullet points or content (3-5 points per slide, each 10-20 words)

Format EXACTLY like this:

SLIDE 1: [Title]
- Point 1
- Point 2
- Point 3

SLIDE 2: [Title]
- Point 1
- Point 2

Continue for all {slides} slides.

First slide should be title slide, last slide should be conclusion/thank you slide.
Make content engaging, professional, and well-structured."""
            
            messages = [{"role": "user", "content": content_prompt}]
            content = llm_handler.get_response(messages, max_tokens=8000, temperature=0.7)
            
            if content.startswith("Error"):
                return f"Failed to generate content: {content}", None
            
            # Create PowerPoint
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Parse slides
            slide_sections = content.split('SLIDE ')
            
            for section in slide_sections:
                section = section.strip()
                if not section or not section[0].isdigit():
                    continue
                
                lines = section.split('\n')
                
                # Extract title
                title_line = lines[0]
                if ':' in title_line:
                    slide_title = title_line.split(':', 1)[1].strip()
                else:
                    slide_title = title_line.strip()
                
                # Extract bullet points
                bullet_points = []
                for line in lines[1:]:
                    line = line.strip()
                    if line.startswith('-') or line.startswith('•'):
                        bullet_points.append(line.lstrip('-•').strip())
                
                # Add slide
                if '1:' in lines[0]:  # Title slide
                    slide = prs.slides.add_slide(prs.slide_layouts[0])
                    title = slide.shapes.title
                    subtitle = slide.placeholders[1]
                    
                    title.text = topic.upper()
                    title.text_frame.paragraphs[0].font.size = Pt(44)
                    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
                    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    
                    subtitle.text = slide_title
                    subtitle.text_frame.paragraphs[0].font.size = Pt(28)
                    subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                else:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    title = slide.shapes.title
                    content = slide.placeholders[1]
                    
                    title.text = slide_title
                    title.text_frame.paragraphs[0].font.size = Pt(32)
                    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
                    
                    text_frame = content.text_frame
                    text_frame.clear()
                    
                    for bullet in bullet_points:
                        p = text_frame.add_paragraph()
                        p.text = bullet
                        p.level = 0
                        p.font.size = Pt(20)
                        p.space_before = Pt(10)
            
            # Save presentation
            safe_topic = "".join(c for c in topic if c.isalnum() or c in (' ', '_')).rstrip()[:50]
            timestamp = time.strftime("%Y%m%d_%H%M%S")
            file_path = self.data_folder / f"{safe_topic.replace(' ', '_')}_{timestamp}.pptx"
            
            prs.save(str(file_path))
            
            Logger.log(f"PowerPoint generated successfully: {file_path}", "PPT")
            return f"I've generated a {slides}-slide presentation on '{topic}' for you, Boss.", str(file_path)
        
        except Exception as e:
            Logger.log(f"Error generating PowerPoint: {e}", "ERROR")
            return f"Failed to generate PowerPoint: {e}", None


# Global instance
ppt_generator = PPTGenerator()